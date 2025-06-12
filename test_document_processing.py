import os
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
from dotenv import load_dotenv

# Add the project root to the Python path
project_root = Path(__file__).parent
sys.path.append(str(project_root))

# Setup logging
from legal_doc_analysis.utils.logging import setup_logging
setup_logging()
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

def extract_text_from_pdf(pdf_path: Path) -> str:
    """Extract text from a PDF file."""
    try:
        import PyPDF2
        text = ""
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {e}")
        return ""

def extract_text_from_docx(docx_path: Path) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document
        doc = Document(docx_path)
        return "\n\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {docx_path}: {e}")
        return ""

def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {pptx_path}: {e}")
        return ""

def extract_text_from_image(image_path: Path) -> str:
    """Extract text from an image using OCR."""
    try:
        import pytesseract
        from PIL import Image
        
        # Set the path to Tesseract executable if not in PATH
        tesseract_cmd = os.getenv('TESSERACT_CMD')
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
            
        img = Image.open(image_path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logger.error(f"Error extracting text from image {image_path}: {e}")
        return ""

def process_document(file_path: Path) -> Dict[str, Any]:
    """Process a document and return its metadata and extracted text."""
    if not file_path.exists():
        return {"error": f"File not found: {file_path}"}
    
    file_info = {
        "file_path": str(file_path),
        "file_name": file_path.name,
        "file_size_mb": file_path.stat().st_size / (1024 * 1024),
        "file_type": file_path.suffix.lower(),
        "extracted_text": ""
    }
    
    try:
        if file_path.suffix.lower() == '.pdf':
            file_info["extracted_text"] = extract_text_from_pdf(file_path)
        elif file_path.suffix.lower() == '.docx':
            file_info["extracted_text"] = extract_text_from_docx(file_path)
        elif file_path.suffix.lower() == '.pptx':
            file_info["extracted_text"] = extract_text_from_pptx(file_path)
        elif file_path.suffix.lower() in ['.jpg', '.jpeg', '.png']:
            file_info["extracted_text"] = extract_text_from_image(file_path)
        else:
            file_info["error"] = f"Unsupported file type: {file_path.suffix}"
            
        return file_info
        
    except Exception as e:
        error_msg = f"Error processing {file_path}: {str(e)}"
        logger.error(error_msg, exc_info=True)
        return {"error": error_msg}

def main():
    # Example usage
    test_files_dir = Path("data/input")
    test_files_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"Looking for test files in: {test_files_dir}")
    
    # List all files in the test directory
    test_files = list(test_files_dir.glob("*"))
    
    if not test_files:
        print(f"No test files found in {test_files_dir}")
        print(f"Please add some test files (PDF, DOCX, PPTX, JPG, PNG) to {test_files_dir}")
        return
    
    print(f"Found {len(test_files)} test file(s):")
    for i, file in enumerate(test_files, 1):
        print(f"{i}. {file.name} ({file.suffix.upper()})")
    
    # Process the first file as an example
    test_file = test_files[0]
    print(f"\nProcessing file: {test_file.name}")
    
    result = process_document(test_file)
    
    print("\nProcessing Results:")
    print(f"File: {result.get('file_name')}")
    print(f"Size: {result.get('file_size_mb'):.2f} MB")
    print(f"Type: {result.get('file_type')}")
    
    if "error" in result:
        print(f"Error: {result['error']}")
    else:
        preview = (result.get('extracted_text', '')[:500] + '...') if result.get('extracted_text') else "No text extracted"
        print(f"Extracted Text (first 500 chars):\n{preview}")

if __name__ == "__main__":
    main()
